"""
文档解析引擎
支持 PDF、Word (.docx)、PPT (.pptx)、Markdown、纯文本
提取文本内容用于 AI 分析和知识索引
"""

from pathlib import Path


def extract_text(filepath):
    """从文件中提取文本内容，返回 (text, metadata)"""
    filepath = Path(filepath)
    suffix = filepath.suffix.lower()

    extractors = {
        '.pdf': _extract_pdf,
        '.docx': _extract_docx,
        '.pptx': _extract_pptx,
        '.md': _extract_text,
        '.txt': _extract_text,
        '.markdown': _extract_text,
    }

    extractor = extractors.get(suffix)
    if not extractor:
        return None, {'error': f'不支持的文件格式: {suffix}'}

    try:
        text, meta = extractor(filepath)
        meta['filename'] = filepath.name
        meta['format'] = suffix
        meta['size_kb'] = round(filepath.stat().st_size / 1024, 1)
        return text, meta
    except Exception as e:
        return None, {'error': f'解析失败: {str(e)}'}


def _extract_pdf(filepath):
    """提取 PDF 文本"""
    import pdfplumber

    text_parts = []
    with pdfplumber.open(filepath) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

    return '\n\n'.join(text_parts), {'pages': page_count}


def _extract_docx(filepath):
    """提取 Word 文档文本"""
    from docx import Document

    doc = Document(filepath)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # 也提取表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)

    return '\n\n'.join(text_parts), {'paragraphs': len(doc.paragraphs)}


def _extract_pptx(filepath):
    """提取 PPT 文本"""
    from pptx import Presentation

    prs = Presentation(filepath)
    text_parts = []
    slide_count = 0

    for slide in prs.slides:
        slide_count += 1
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text.strip())
        if slide_texts:
            text_parts.append(f"[Slide {slide_count}]\n" + '\n'.join(slide_texts))

    return '\n\n'.join(text_parts), {'slides': slide_count}


def _extract_text(filepath):
    """提取纯文本/Markdown"""
    text = filepath.read_text(encoding='utf-8')
    lines = len(text.split('\n'))
    return text, {'lines': lines}


def generate_summary_prompt(text, filename):
    """生成用于 AI 摘要的 prompt"""
    # 截断过长的文本（保留前 3000 字）
    truncated = text[:3000] if len(text) > 3000 else text
    was_truncated = len(text) > 3000

    prompt = f"""请对以下文档内容生成一份结构化摘要（用于知识管理）。

文件名：{filename}
{'（内容已截断，仅展示前3000字）' if was_truncated else ''}

---
{truncated}
---

请输出以下格式的 JSON（不要加 markdown 代码块标记）：
{{
  "title": "文档标题或主题",
  "category": "方法论/行业分析/SOP流程/学习笔记/周报/其他",
  "key_points": ["核心要点1", "核心要点2", "核心要点3"],
  "keywords": ["关键词1", "关键词2", "关键词3", "关键词4", "关键词5"],
  "one_line_summary": "一句话概括这个文档的核心内容",
  "relevance_to_ops": "这个文档对策略运营/产品运营工作的价值和应用场景"
}}"""

    return prompt
